import copy
import six
from pptx import Presentation
import boto3
import io
import time

def delete_slide(presentation,  index):
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
def load_presentation(name):
    s3 = boto3.client('s3')
    s3_response_object = s3.get_object(Bucket='userpptx', Key=name)
    object_content = s3_response_object['Body'].read()

    prs = Presentation(io.BytesIO(object_content))
    print("LOADED PRESENTATION...")
    return prs

def upload_presentation(pres, name):
    pres.save("/tmp/test_presentating.pptx")
    s3 = boto3.client('s3')
    print("UPLOADING PRESENTATION...")
    s3.upload_file("/tmp/test_presentating.pptx", "userpptx", name)
    url = s3.generate_presigned_url(
            ClientMethod='get_object',
            Params={'Bucket': 'userpptx', 'Key': name},
            ExpiresIn=3600  # URL valid for 1 hour
        )
    return url

def move_slide(presentation, old_index, new_index):
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def edit_text(shape, text):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            run.text = text
            run.font.size = font.size
            run.font.bold = font.bold
            run.font.italic = font.italic
            run.font.underline = font.underline
            run.font.color.rgb = font.color.rgb

def _duplicate_slide(pres, place, title="Title", summary="Summary", index=3):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[10]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )
    if template.background.fill.type:
        copied_slide.background.fill.solid()
        copied_slide.background.fill.fore_color.rgb = template.background.fill.fore_color.rgb
    print("Duplicating slide...")
    edit_text(copied_slide.shapes[-1], title)
    edit_text(copied_slide.shapes[-2], summary)
    move_slide(pres, len(pres.slides) - 1, place)
    print("Created slide...")
    return "SUCCESS!"

def _make_presentation(titles, descriptions):
    print("Making presentation...")
    pres = load_presentation("Design 2 MAM-BOOK.pptx")

    titles = titles.split("^")
    descriptions = descriptions.split("^")
    presentation_values = dict(zip(titles, descriptions))
    for i, (title, description) in enumerate(presentation_values.items()):
        title = title.replace("#", "").replace("\n", "").replace("/", "")
        description = description.replace("#", "").replace("\n", "").replace("/", "")
        print("Creating slides...")
        print("Title:", title)
        print("Description:", description)
        status = _duplicate_slide(pres=pres, title=title, summary=description, place=(i+3))
        # delete_slide(pres, -3)
        # delete_slide(pres, -2)
        # delete_slide(pres, 2)
        if status=="SUCCESS!":
            print("Slide created successfully...")
    try :
        return upload_presentation(pres, "test" + time.strftime("%Y-%m-%d-%H-%M-%S") + ".pptx")
    except Exception as e:
        print(e)
        return "ERROR"