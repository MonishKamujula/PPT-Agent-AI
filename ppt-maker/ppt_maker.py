import copy
import six
from pptx import Presentation

def _duplicate_slide(pres, title="Title", summary="Summary", index=3):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[10]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )
    if template.background.fill.type:
        copied_slide.background.fill.solid()
        copied_slide.background.fill.fore_color.rgb = template.background.fill.fore_color.rgb
    print("Duplicating slide...")
    copied_slide.shapes[-1].text = title
    copied_slide.shapes[-2].text = summary
    pres.save("C:\\Monish\\deep-learning\\MAM-book-summary-presenter\\ppt-maker\\pptx-presentations\\modified.pptx")
    print("Created slide...")
    return "SUCCESS!"

def _make_presentation(titles, descriptions):
    print("Making presentation...")
    pres = Presentation("C:\\Monish\\deep-learning\\MAM-book-summary-presenter\\ppt-maker\\pptx-presentations\\Design 2 MAM-BOOK.pptx")

    titles = titles.split("^")
    descriptions = descriptions.split("^")
    presentation_values = dict(zip(titles, descriptions))
    for title, description in presentation_values.items():
        title = title.replace("#", "").replace("\n", "").replace("/", "")
        description = description.replace("#", "").replace("\n", "").replace("/", "")
        print("Creating slides...")
        print("Title:", title)
        print("Description:", description)
        status = _duplicate_slide(pres, title, description)
        if status:
            print("Slide created successfully...")
    return "SUCCESS!"

